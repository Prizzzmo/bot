
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests
from io import BytesIO
from PIL import Image
import markdown
import html2text
import time

def md_to_plaintext(md_text):
    """Конвертирует маркдаун в простой текст"""
    html = markdown.markdown(md_text)
    h = html2text.HTML2Text()
    h.ignore_links = True
    return h.handle(html).strip()

def download_and_verify_image(url, file_path):
    """
    Скачивает изображение, проверяет его валидность и возвращает путь к файлу
    или None в случае ошибки
    """
    try:
        # Если файл уже существует, пробуем сразу его проверить
        if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
            try:
                img = Image.open(file_path)
                img.verify()  # Проверка, что файл является изображением
                return file_path
            except Exception:
                print(f"Существующий файл {file_path} поврежден. Пробуем скачать заново.")
                # Если проверка не удалась, продолжаем и пробуем скачать заново
                pass
        
        # Скачиваем изображение
        response = requests.get(url, timeout=10)
        if response.status_code != 200:
            print(f"Ошибка при скачивании изображения: HTTP {response.status_code}")
            return None
            
        # Проверяем скачанное изображение
        try:
            img = Image.open(BytesIO(response.content))
            img.verify()  # Проверка, что данные являются изображением
            
            # Сохраняем валидное изображение
            with open(file_path, 'wb') as f:
                f.write(response.content)
            
            return file_path
        except Exception as e:
            print(f"Скачанное изображение не является валидным: {e}")
            return None
            
    except Exception as e:
        print(f"Ошибка при скачивании/проверке изображения {url}: {e}")
        return None

def create_presentation(md_file_path, output_pptx_path):
    """
    Преобразует Markdown-файл в презентацию PowerPoint.
    
    Args:
        md_file_path (str): Путь к файлу Markdown
        output_pptx_path (str): Путь для сохранения презентации PowerPoint
    """
    # Создаем папку для временных изображений, если она не существует
    images_dir = 'presentation_images'
    os.makedirs(images_dir, exist_ok=True)
    
    # Создаем новую презентацию
    prs = Presentation()
    
    # Читаем содержимое Markdown-файла
    with open(md_file_path, 'r', encoding='utf-8') as file:
        md_content = file.read()
    
    # Разбиваем контент на секции по заголовкам первого уровня
    sections = re.split(r'## ', md_content)
    
    # Получаем заголовок всей презентации
    title_match = re.match(r'# (.*?)\n', sections[0])
    main_title = title_match.group(1) if title_match else "Презентация проекта"
    
    # Создаем титульный слайд
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = main_title
    subtitle.text = "Проектная презентация"
    
    # Добавляем изображение для титульного слайда
    image_path = os.path.join(images_dir, "russia_history.jpg")
    image_url = "https://histrf.ru/uploads/media/default/0001/02/0e4f4e9d11f6bc76c7c69de95b3ab6f3c6fb7d3b.jpeg"
    
    verified_image = download_and_verify_image(image_url, image_path)
    if verified_image:
        try:
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            slide.shapes.add_picture(verified_image, left, top, width=width)
            print(f"Изображение для титульного слайда успешно добавлено")
        except Exception as e:
            print(f"Ошибка при добавлении изображения на титульный слайд: {e}")
    
    # Создаем слайд с содержанием
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Содержание"
    
    # Извлекаем содержание из первой секции
    content_match = re.search(r'### Содержание\n(.*?)(?=\n##|\Z)', sections[0], re.DOTALL)
    if content_match:
        content_text = content_match.group(1)
        # Удаляем нумерацию и ссылки
        content_text = re.sub(r'\d+\.\s*\[(.*?)\].*?', r'\1', content_text)
        content.text = content_text
    
    # Обрабатываем каждую секцию (кроме первой, которая содержит заголовок и введение)
    for i, section in enumerate(sections[1:], 1):
        # Получаем заголовок секции
        section_title = section.split('\n', 1)[0].strip()
        
        # Создаем слайд для заголовка секции
        section_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = section_title
        
        # Добавляем релевантное изображение в зависимости от заголовка
        image_url = None
        if "Архитектура" in section_title:
            image_url = "https://miro.medium.com/v2/resize:fit:1200/0*9pGRpgfmnqPEBuFj.png"
            image_name = "архитектура_системы.jpg"
        elif "ИИ" in section_title or "Gemini" in section_title:
            image_url = "https://storage.googleapis.com/gweb-uniblog-publish-prod/images/Gemini_Features_2.max-1000x1000.png"
            image_name = "интеграция_с_ии.jpg"
        elif "Компонент" in section_title:
            image_url = "https://cdn-icons-png.flaticon.com/512/2530/2530302.png"
            image_name = f"компоненты_системы.jpg"
        elif "Карт" in section_title or "Визуализаци" in section_title:
            image_url = "https://cdn-icons-png.flaticon.com/512/854/854878.png"
            image_name = "карты_и_визуализации.jpg"
        elif "Тестирован" in section_title:
            image_url = "https://cdn-icons-png.flaticon.com/512/2620/2620343.png"
            image_name = "тестирование.jpg"
        elif "Аналитич" in section_title:
            image_url = "https://cdn-icons-png.flaticon.com/512/1925/1925059.png"
            image_name = "аналитическая_система.jpg"
        else:
            # Для других разделов используем стандартное изображение
            image_url = "https://cdn-icons-png.flaticon.com/512/4481/4481249.png" 
            image_name = f"{section_title.lower().replace(' ', '_').replace(':', '_')}.jpg"
        
        if image_url:
            image_path = os.path.join(images_dir, image_name)
            verified_image = download_and_verify_image(image_url, image_path)
            
            if verified_image:
                try:
                    left = Inches(3)
                    top = Inches(2.5)
                    width = Inches(4)
                    slide.shapes.add_picture(verified_image, left, top, width=width)
                    print(f"Изображение для секции '{section_title}' успешно добавлено")
                except Exception as e:
                    print(f"Ошибка при добавлении изображения для секции '{section_title}': {e}")
        
        # Разбиваем секцию на подсекции и создаем слайды для каждой подсекции
        subsections = re.split(r'### ', section)[1:] if '### ' in section else []
        
        # Если нет подсекций, создаем слайд с основным содержанием
        if not subsections:
            main_content = section.split('\n', 1)[1] if '\n' in section else ""
            
            # Убираем Markdown форматирование для простого текста
            main_content = md_to_plaintext(main_content)
            
            # Создаем слайд с содержанием
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = section_title
            
            # Ограничиваем размер текста, если он слишком длинный
            if len(main_content) > 500:
                main_content = main_content[:500] + "..."
            
            content.text = main_content
        
        # Обрабатываем каждую подсекцию
        for subsection in subsections:
            # Получаем заголовок подсекции
            subsection_title = subsection.split('\n', 1)[0].strip()
            
            # Создаем слайд для подсекции
            subsection_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(subsection_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = subsection_title
            
            # Получаем содержание подсекции
            subsection_content = subsection.split('\n', 1)[1] if '\n' in subsection else ""
            
            # Преобразуем маркдаун в текст
            subsection_content = md_to_plaintext(subsection_content)
            
            # Ограничиваем размер текста
            if len(subsection_content) > 500:
                subsection_content = subsection_content[:500] + "..."
            
            content.text = subsection_content
    
    # Заключительный слайд
    final_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(final_slide_layout)
    title = slide.shapes.title
    title.text = "Спасибо за внимание!"
    
    # Добавляем текст с контактной информацией
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textframe = textbox.text_frame
    textframe.text = "© 2025 Образовательный бот по истории России"
    textframe.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Добавляем изображение флага России
    image_path = os.path.join(images_dir, "russia_flag.jpg")
    image_url = "https://upload.wikimedia.org/wikipedia/commons/thumb/f/f3/Flag_of_Russia.svg/800px-Flag_of_Russia.svg.png"
    
    verified_image = download_and_verify_image(image_url, image_path)
    if verified_image:
        try:
            left = Inches(3)
            top = Inches(4)
            width = Inches(4)
            slide.shapes.add_picture(verified_image, left, top, width=width)
            print(f"Изображение для заключительного слайда успешно добавлено")
        except Exception as e:
            print(f"Ошибка при добавлении изображения на заключительный слайд: {e}")
    
    # Сохраняем презентацию
    prs.save(output_pptx_path)
    print(f"Презентация PowerPoint успешно создана и сохранена: {output_pptx_path}")
    return output_pptx_path

if __name__ == "__main__":
    # Создаем презентацию PowerPoint из Markdown файла
    create_presentation('project_presentation.md', 'История_России_образовательный_бот_презентация.pptx')
